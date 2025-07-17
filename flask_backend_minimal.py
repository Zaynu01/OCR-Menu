from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import os
import tempfile
from werkzeug.utils import secure_filename
import pandas as pd
from pptx import Presentation

# Import your existing helper functions
from utils.helper import (
    PowerPointHelper,
    clone_template_slide,
    remove_placeholders_from_cloned_slides,
)

app = Flask(__name__)
CORS(app)

# Configuration
UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'pptx', 'xlsx'}
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max

os.makedirs(UPLOAD_FOLDER, exist_ok=True)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def process_powerpoint_files(ppt_file_path, excel_file_path):
    """Your existing processing logic wrapped in a function"""
    results = []
    
    try:
        # Load Excel
        df = pd.read_excel(excel_file_path, engine="openpyxl")
        
        # Validate columns
        expected_columns = ["old_path", "reference_old", "new_path", "reference_new", "commentaire", "conclusion"]
        cols_lower = [col.lower() for col in df.columns[:6]]
        if df.empty or not all(col.lower() in cols_lower for col in expected_columns):
            return False, "Excel must contain required columns.", None, []
        
        # Filter rows
        df_filtered = df.dropna(subset=["old_path", "new_path"])
        df_filtered = df_filtered[(df_filtered["old_path"].str.strip() != "") & (df_filtered["new_path"].str.strip() != "")]
        df_filtered = df_filtered.reset_index(drop=True)
        
        if df_filtered.empty:
            return False, "No valid rows found.", None, []
        
        # Extract columns
        old_paths = df_filtered["old_path"].tolist()
        reference_old = df_filtered["reference_old"].fillna("").tolist()
        new_paths = df_filtered["new_path"].tolist()
        reference_new = df_filtered["reference_new"].fillna("").tolist()
        commentaire = df_filtered["commentaire"].fillna("").tolist()
        conclusion = df_filtered["conclusion"].fillna("").tolist()
        
        # Create working copy
        temp_ppt = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        ppt_path = temp_ppt.name
        temp_ppt.close()
        
        with open(ppt_file_path, "rb") as src:
            with open(ppt_path, "wb") as dst:
                dst.write(src.read())
        
        # Clone slides
        num_rows = len(df_filtered)
        ppt_path = clone_template_slide(ppt_path, num_rows)
        
        # Validate slides
        presentation = Presentation(ppt_path)
        if len(presentation.slides) <= 1:
            return False, "PowerPoint must have at least 2 slides.", None, []
        
        # Process images and text
        image_tool = PowerPointHelper()
        
        for i in range(num_rows):
            slide_number = i + 2
            left_img_path = old_paths[i]
            right_img_path = new_paths[i]
            
            # Replace images
            if os.path.exists(left_img_path):
                output_path = image_tool.replace_left_image(ppt_path, slide_number, left_img_path)
                if output_path:
                    ppt_path = output_path
                    results.append(f"✅ Slide {slide_number}: Left image replaced")
                else:
                    results.append(f"❌ Slide {slide_number}: Failed to replace left image")
            else:
                results.append(f"❌ Slide {slide_number}: Left image not found")
            
            if os.path.exists(right_img_path):
                output_path = image_tool.replace_right_image(ppt_path, slide_number, right_img_path)
                if output_path:
                    ppt_path = output_path
                    results.append(f"✅ Slide {slide_number}: Right image replaced")
                else:
                    results.append(f"❌ Slide {slide_number}: Failed to replace right image")
            else:
                results.append(f"❌ Slide {slide_number}: Right image not found")
            
            # Add text if both images exist
            if os.path.exists(left_img_path) and os.path.exists(right_img_path):
                text_tool = PowerPointHelper()
                ppt_path = text_tool.add_texts_from_excel(
                    ppt_path,
                    commentaires=[commentaire[i]],
                    conclusions=[conclusion[i]],
                    reference_old=[reference_old[i]],
                    reference_new=[reference_new[i]]
                )
            else:
                results.append(f"⚠️ Slide {slide_number}: Skipped text replacement")
        
        # Remove placeholders
        ppt_path = remove_placeholders_from_cloned_slides(ppt_path, start_index=2)
        
        return True, "Processing completed successfully!", ppt_path, results
        
    except Exception as e:
        return False, f"Error: {str(e)}", None, results

@app.route('/api/upload', methods=['POST'])
def upload_files():
    print('request.files:', request.files)
    print('request.form:', request.form)
    print('request.content_type:', request.content_type)
    try:
        if 'ppt_file' not in request.files or 'excel_file' not in request.files:
            return jsonify({'error': 'Both files required'}), 400
        
        ppt_file = request.files['ppt_file']
        excel_file = request.files['excel_file']
        
        if not (ppt_file.filename and excel_file.filename):
            return jsonify({'error': 'No files selected'}), 400
        
        if not (allowed_file(ppt_file.filename) and allowed_file(excel_file.filename)):
            return jsonify({'error': 'Invalid file types'}), 400
        
        # Save files
        ppt_filename = secure_filename(ppt_file.filename)
        excel_filename = secure_filename(excel_file.filename)
        
        ppt_path = os.path.join(app.config['UPLOAD_FOLDER'], ppt_filename)
        excel_path = os.path.join(app.config['UPLOAD_FOLDER'], excel_filename)
        
        ppt_file.save(ppt_path)
        excel_file.save(excel_path)
        
        # Process files
        success, message, result_path, results = process_powerpoint_files(ppt_path, excel_path)
        
        # Clean up
        os.remove(ppt_path)
        os.remove(excel_path)
        
        if success:
            return jsonify({
                'success': True,
                'message': message,
                'results': results,
                'download_token': os.path.basename(result_path)
            })
        else:
            return jsonify({
                'success': False,
                'message': message,
                'results': results
            }), 400
            
    except Exception as e:
        return jsonify({'error': f'Server error: {str(e)}'}), 500

@app.route('/api/download/<token>')
def download_file(token):
    try:
        file_path = os.path.join(tempfile.gettempdir(), token)
        if os.path.exists(file_path):
            return send_file(file_path, as_attachment=True, download_name='modified_presentation.pptx')
        else:
            return jsonify({'error': 'File not found'}), 404
    except Exception as e:
        return jsonify({'error': f'Download error: {str(e)}'}), 500

@app.route('/api/health')
def health_check():
    return jsonify({'status': 'healthy'})

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)
